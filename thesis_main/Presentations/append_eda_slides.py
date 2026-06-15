"""Append 3 EDA slides to the existing defense pptx, preserving user edits.

Insert them after slide 14 (Geospatial Features Overview) by manual drag-reorder
in PowerPoint/Canva. We append at end to avoid disturbing existing slide order.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent.parent
SRC_PPTX = ROOT / "thesis_main/Presentations/defense_2026-05-09.pptx"
OUT_PPTX = ROOT / "thesis_main/Presentations/defense_2026-05-09_with_eda.pptx"

EDA = ROOT / "EDA"
QGIS = ROOT / "thesis_main/QGIS/outputs"

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
AMBER = RGBColor(0xD4, 0xA8, 0x43)
SLATE = RGBColor(0xEE, 0xF1, 0xF6)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x5F, 0x6D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_header(slide, slide_w, title_text, eyebrow_text):
    # Top navy strip
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(0.55)
    )
    strip.line.fill.background()
    strip.fill.solid()
    strip.fill.fore_color.rgb = NAVY
    tf = strip.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.12)
    tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = eyebrow_text
    p.font.name = "Inter"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = AMBER

    # Title row
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.7), slide_w - Inches(1.0), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Inter"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Amber underline accent
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.45), Inches(0.9), Inches(0.05)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMBER


def add_image_fit(slide, image_path, left, top, max_w, max_h):
    """Add image scaled to fit max_w/max_h while preserving aspect ratio."""
    pic = slide.shapes.add_picture(str(image_path), left, top)
    nw, nh = pic.width, pic.height
    rw = max_w / nw
    rh = max_h / nh
    r = min(rw, rh)
    pic.width = int(nw * r)
    pic.height = int(nh * r)
    # center within bounding box
    pic.left = left + (max_w - pic.width) // 2
    pic.top = top + (max_h - pic.height) // 2
    return pic


def add_bullets(slide, left, top, w, h, items):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + item
        p.font.name = "Inter"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)


def add_speaker_line(slide, left, top, w, h, text):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = SLATE
    # amber left border
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), h)
    border.line.fill.background()
    border.fill.solid()
    border.fill.fore_color.rgb = AMBER

    tf = box.text_frame
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Inter"
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = NAVY


def set_speaker_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def add_blank_slide(prs):
    # pptxgenjs-built deck only has 1 layout ('DEFAULT'); use it.
    return prs.slides.add_slide(prs.slide_layouts[0])


def main():
    prs = Presentation(str(SRC_PPTX))
    sw = prs.slide_width
    sh = prs.slide_height
    print(f"Slide size: {sw/914400:.2f} x {sh/914400:.2f} in")

    # ---- EDA Slide A: Geographic & Market Coverage ----
    s = add_blank_slide(prs)
    add_header(s, sw, "Geographic & Market Coverage",
               "EDA  /  DATA STRATEGY")

    # Hero map (left, large)
    map_left = Inches(0.5)
    map_top = Inches(1.7)
    map_w = Inches(7.5)
    map_h = Inches(4.3)
    add_image_fit(s, QGIS / "osm-map_properties.png",
                  map_left, map_top, map_w, map_h)

    # Caption under map
    cap = s.shapes.add_textbox(map_left, map_top + map_h, map_w, Inches(0.3))
    p = cap.text_frame.paragraphs[0]
    p.text = "Open-market properties color-coded by LGU. Mactan Island visibly separated from mainland Cebu by water."
    p.font.name = "Inter"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = GRAY

    # Right inset: city/segment heatmap
    inset_left = Inches(8.2)
    inset_top = Inches(1.7)
    inset_w = Inches(4.8)
    inset_h = Inches(2.8)
    add_image_fit(s, EDA / "city_segment_heatmap.png",
                  inset_left, inset_top, inset_w, inset_h)

    # Right bullets
    add_bullets(s, inset_left, inset_top + inset_h + Inches(0.15),
                inset_w, Inches(2.0), [
        "Lapu-Lapu (559) + Cebu City (508) anchor open-market sample",
        "Talisay (85), Minglanilla (71) — smaller samples, wider error",
        "Mactan separation motivates osmnx network distance over Haversine",
    ])

    # Speaker line at bottom
    add_speaker_line(s, Inches(0.5), Inches(6.7),
                     sw - Inches(1.0), Inches(0.55),
                     "Geographic spread foreshadows the per-LGU MAPE story — sample size and "
                     "physical separation explain most of the variance in error.")

    set_speaker_notes(s,
        "Bullet talking points: We collected 1,619 open-market rows across 6 LGUs. "
        "Lapu-Lapu and Cebu City dominate, giving us our most credible single-LGU benchmarks. "
        "Mactan island's water boundary is the key reason we use osmnx network distance: "
        "Haversine would treat it as 5 km from CBP when actual road distance via the Mactan "
        "bridge is closer to 12-15 km. The smaller-sample LGUs — Talisay 85, Minglanilla 71 — "
        "are flagged as indicative in our results section."
    )

    # ---- EDA Slide B: Price Distribution by City ----
    s = add_blank_slide(prs)
    add_header(s, sw, "Price Distribution by City",
               "EDA  /  TARGET DISTRIBUTION")

    # Hero boxplot (centered, large)
    img_left = Inches(0.5)
    img_top = Inches(1.7)
    img_w = Inches(8.5)
    img_h = Inches(4.8)
    add_image_fit(s, EDA / "price_by_city_open_market.png",
                  img_left, img_top, img_w, img_h)

    # Right column bullets
    txt_left = Inches(9.2)
    txt_top = Inches(1.8)
    txt_w = Inches(3.8)
    txt_h = Inches(4.5)
    add_bullets(s, txt_left, txt_top, txt_w, txt_h, [
        "Cebu City: median ≈ PHP 150K/sqm, widest spread",
        "Mandaue + Lapu-Lapu: median ≈ PHP 115–125K/sqm",
        "Consolacion + Talisay + Minglanilla: median ≈ PHP 55–75K/sqm",
        "Heavy right-skew → log-transform for OLS baseline",
        "Tier separation justifies city dummies as primary structural feature",
    ])

    add_speaker_line(s, Inches(0.5), Inches(6.7),
                     sw - Inches(1.0), Inches(0.55),
                     "Three price tiers across Metro Cebu — and the boundary between them is "
                     "exactly what the model has to learn. The wide IQR in Cebu City foreshadows "
                     "the residual error in the luxury segment.")

    set_speaker_notes(s,
        "The boxplot shows a clear three-tier structure: Cebu City highest, Mandaue/Lapu-Lapu in "
        "the middle, the smaller LGUs lowest. Cebu City's IQR is widest because it spans both "
        "mid-tier subdivisions and luxury developments — that mix-tier dispersion is part of "
        "why the residual error is highest in that segment. The right-skew of the raw distribution "
        "(visible as outliers above each box) is what motivated the log-transform for OLS — "
        "without it OLS R-squared was -45 on the test set; with it, OLS reaches 0.394."
    )

    # ---- EDA Slide C: Feature Diagnostics ----
    s = add_blank_slide(prs)
    add_header(s, sw, "Feature Diagnostics: Correlations",
               "EDA  /  FEATURE SELECTION RATIONALE")

    # Two images side-by-side
    left_img = Inches(0.5)
    right_img = Inches(6.9)
    img_top = Inches(1.7)
    img_w = Inches(6.1)
    img_h = Inches(4.5)

    add_image_fit(s, EDA / "cbd_distance_corr.png",
                  left_img, img_top, img_w, img_h)
    add_image_fit(s, EDA / "target_corr_open_market.png",
                  right_img, img_top, img_w, img_h)

    # Sub-captions under each
    cap1 = s.shapes.add_textbox(left_img, img_top + img_h + Inches(0.05),
                                img_w, Inches(0.3))
    p = cap1.text_frame.paragraphs[0]
    p.text = "CBD distance correlations — retained nodes after multicollinearity check"
    p.font.name = "Inter"; p.font.size = Pt(10); p.font.italic = True
    p.font.color.rgb = GRAY

    cap2 = s.shapes.add_textbox(right_img, img_top + img_h + Inches(0.05),
                                img_w, Inches(0.3))
    p = cap2.text_frame.paragraphs[0]
    p.text = "Spearman correlations of features with price_per_sqm (open_market)"
    p.font.name = "Inter"; p.font.size = Pt(10); p.font.italic = True
    p.font.color.rgb = GRAY

    add_speaker_line(s, Inches(0.5), Inches(6.65),
                     sw - Inches(1.0), Inches(0.6),
                     "Two diagnostics drove feature selection: the CBD correlation matrix "
                     "(IT Park dropped at r=0.99 with CBP) and the target-correlation ranking "
                     "(MCRAI recreation, BIR zonal medians, MCRAI composite as top positives).")

    set_speaker_notes(s,
        "Left panel: CBD distance multicollinearity check. The 8 retained nodes still show high "
        "intra-cluster correlations — Talisay-Tabunok and SRP at 0.97, Naga and Talisay-Tabunok at "
        "0.98 — but these are kept because tree-based models handle correlated features without the "
        "unstable coefficients that would break OLS. The IT Park drop happened earlier at r=0.99 "
        "with Cebu Business Park. Right panel: Spearman correlations with price_per_sqm. MCRAI "
        "recreation tops the positive list — supports including it in the composite. BIR zonal "
        "medians correlate strongly but are excluded as features (used only for valuation gap "
        "analysis). Floor area shows strongest negative correlation — reflects per-sqm pricing "
        "behavior in larger developments, not a true price-floor relationship."
    )

    prs.save(str(OUT_PPTX))
    print(f"Saved: {OUT_PPTX.name}")
    print(f"Total slides now: {len(prs.slides)}")


if __name__ == "__main__":
    main()
