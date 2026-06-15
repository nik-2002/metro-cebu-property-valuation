"""Insert 5 new slides into final_Draft_Estreba_thesis_proposal.pptx in place.

Insertions (final positions):
  27: Model Performance Visualized (NEW chart slide)
  29: RF Actual vs Predicted - with image (NEW)
  34: Per-LGU MAPE Visualized (NEW chart slide)
  41: Research Questions Answered (NEW)
  43: Recommendations (NEW)

Charts are rendered with matplotlib and embedded as images so they match
the navy / amber design system.
"""
from pathlib import Path
import tempfile

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent.parent
SRC = ROOT / "final_Draft_Estreba_thesis_proposal.pptx"
EDA = ROOT / "EDA"

# Theme
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
AMBER = RGBColor(0xD4, 0xA8, 0x43)
SLATE = RGBColor(0xEE, 0xF1, 0xF6)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x5F, 0x6D)

NAVY_HEX = "#1B2A4A"
AMBER_HEX = "#D4A843"
SLATE_HEX = "#EEF1F6"
GRAY_HEX = "#555F6D"

TMP = Path(tempfile.gettempdir()) / "thesis_charts"
TMP.mkdir(exist_ok=True)


# ----------------------- chart helpers -----------------------
def render_model_comparison_chart():
    out = TMP / "model_comparison_chart.png"
    models = ["OLS\nHedonic", "Random Forest\n(deployed)", "XGBoost\n(tuned)"]
    r2 = [0.083, 0.807, 0.557]
    mape = [201.6, 59.28, 58.93]
    mae = [9.82, 4.95, 6.06]

    fig, axes = plt.subplots(1, 3, figsize=(13, 4.2), dpi=150)
    fig.patch.set_facecolor("white")
    colors = [GRAY_HEX, AMBER_HEX, NAVY_HEX]

    # R²
    ax = axes[0]
    bars = ax.bar(models, r2, color=colors, edgecolor="white")
    ax.set_title("R²  (higher = better)", fontsize=13, color=NAVY_HEX, fontweight="bold")
    ax.set_ylim(0, 1.0)
    ax.spines[["top", "right"]].set_visible(False)
    for bar, v in zip(bars, r2):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.02, f"{v:.3f}",
                ha="center", fontsize=11, color=NAVY_HEX, fontweight="bold")

    # MAPE (clip OLS visually with annotation)
    ax = axes[1]
    mape_clip = [min(v, 100) for v in mape]
    bars = ax.bar(models, mape_clip, color=colors, edgecolor="white")
    ax.set_title("MAPE %  (lower = better)", fontsize=13, color=NAVY_HEX, fontweight="bold")
    ax.set_ylim(0, 110)
    ax.spines[["top", "right"]].set_visible(False)
    for bar, v in zip(bars, mape):
        label = f"{v:.1f}%" + (" ↑" if v > 100 else "")
        ax.text(bar.get_x() + bar.get_width() / 2,
                min(v, 100) + 2, label,
                ha="center", fontsize=11, color=NAVY_HEX, fontweight="bold")
    ax.text(0.0, 105, "OLS clipped at 100%; actual = 201.6%",
            fontsize=9, color=GRAY_HEX, style="italic")

    # MAE PHP M
    ax = axes[2]
    bars = ax.bar(models, mae, color=colors, edgecolor="white")
    ax.set_title("MAE  (PHP M, lower = better)", fontsize=13, color=NAVY_HEX, fontweight="bold")
    ax.set_ylim(0, 12)
    ax.spines[["top", "right"]].set_visible(False)
    for bar, v in zip(bars, mae):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.2, f"{v:.2f}",
                ha="center", fontsize=11, color=NAVY_HEX, fontweight="bold")

    for ax in axes:
        for lbl in ax.get_xticklabels():
            lbl.set_fontsize(10)
            lbl.set_color(NAVY_HEX)
        for lbl in ax.get_yticklabels():
            lbl.set_color(GRAY_HEX)
        ax.tick_params(colors=GRAY_HEX)

    plt.tight_layout()
    plt.savefig(out, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def render_per_lgu_mape_chart():
    out = TMP / "per_lgu_mape_chart.png"
    rows = [
        ("Consolacion", 21, 33.3),
        ("Lapu-Lapu", 108, 40.0),
        ("Cebu City", 86, 42.2),
        ("Minglanilla", 12, 47.2),
        ("Mandaue", 64, 51.4),
        ("Talisay", 13, 65.8),
    ]
    rows = sorted(rows, key=lambda r: r[2])
    labels = [f"{r[0]}  (n={r[1]})" for r in rows]
    values = [r[2] for r in rows]

    fig, ax = plt.subplots(figsize=(11, 4.8), dpi=150)
    fig.patch.set_facecolor("white")
    colors = []
    for r in rows:
        if r[0] == "Lapu-Lapu":
            colors.append(AMBER_HEX)  # most credible
        elif r[1] < 20:
            colors.append("#B5BCC8")  # caution: small sample
        else:
            colors.append(NAVY_HEX)

    y = list(range(len(rows)))
    bars = ax.barh(y, values, color=colors, edgecolor="white")
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=11, color=NAVY_HEX)
    ax.invert_yaxis()
    ax.set_xlabel("MAPE %", fontsize=11, color=GRAY_HEX)
    ax.set_xlim(0, 75)
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(colors=GRAY_HEX)
    for lbl in ax.get_xticklabels():
        lbl.set_color(GRAY_HEX)

    for bar, v in zip(bars, values):
        ax.text(v + 1.0, bar.get_y() + bar.get_height() / 2,
                f"{v:.1f}%", va="center", fontsize=11, color=NAVY_HEX, fontweight="bold")

    legend_handles = [
        mpatches.Patch(color=AMBER_HEX, label="Most credible (largest sample)"),
        mpatches.Patch(color=NAVY_HEX, label="Adequate sample"),
        mpatches.Patch(color="#B5BCC8", label="Small sample (n<20) — indicative"),
    ]
    ax.legend(handles=legend_handles, loc="lower right", fontsize=9, frameon=False)
    ax.set_title("MAPE by LGU (test set, sorted ascending)",
                 fontsize=13, color=NAVY_HEX, fontweight="bold", loc="left")

    plt.tight_layout()
    plt.savefig(out, dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


# ----------------------- slide helpers -----------------------
def add_header(slide, slide_w, eyebrow, title):
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(0.55))
    strip.line.fill.background()
    strip.fill.solid(); strip.fill.fore_color.rgb = NAVY
    tf = strip.text_frame
    tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = eyebrow
    p.font.name = "Inter"; p.font.size = Pt(11); p.font.bold = True
    p.font.color.rgb = AMBER

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7),
                                         slide_w - Inches(1.0), Inches(0.7))
    tf = title_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Inter"; p.font.size = Pt(28); p.font.bold = True
    p.font.color.rgb = NAVY

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(1.45), Inches(0.9), Inches(0.05))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = AMBER


def add_image_fit(slide, image_path, left, top, max_w, max_h):
    pic = slide.shapes.add_picture(str(image_path), left, top)
    nw, nh = pic.width, pic.height
    r = min(max_w / nw, max_h / nh)
    pic.width = int(nw * r); pic.height = int(nh * r)
    pic.left = left + (max_w - pic.width) // 2
    pic.top = top + (max_h - pic.height) // 2
    return pic


def add_speaker_strip(slide, slide_w, text):
    left = Inches(0.5); top = Inches(6.7)
    w = slide_w - Inches(1.0); h = Inches(0.55)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    box.line.fill.background()
    box.fill.solid(); box.fill.fore_color.rgb = SLATE
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), h)
    border.line.fill.background()
    border.fill.solid(); border.fill.fore_color.rgb = AMBER
    tf = box.text_frame
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Inter"; p.font.size = Pt(11); p.font.italic = True
    p.font.color.rgb = NAVY


def add_card(slide, left, top, w, h, eyebrow, title, body, accent=AMBER):
    # background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    bg.line.color.rgb = RGBColor(0xDD, 0xE2, 0xEA)
    bg.line.width = Pt(0.75)
    bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), h)
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = accent

    # text content
    pad = Inches(0.25)
    tb = slide.shapes.add_textbox(left + pad, top + Inches(0.15),
                                  w - pad - Inches(0.2), h - Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = eyebrow
    p.font.name = "Inter"; p.font.size = Pt(10); p.font.bold = True
    p.font.color.rgb = accent

    p = tf.add_paragraph(); p.text = title
    p.font.name = "Inter"; p.font.size = Pt(15); p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_before = Pt(2); p.space_after = Pt(4)

    p = tf.add_paragraph(); p.text = body
    p.font.name = "Inter"; p.font.size = Pt(11); p.font.color.rgb = DARK


def set_notes(slide, text):
    try:
        tf = slide.notes_slide.notes_text_frame
        if tf is None:
            return
        tf.text = text
    except Exception:
        # Notes placeholder not available on this layout; skip silently.
        pass


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[0])


# ----------------------- slide builders -----------------------
def build_model_perf_slide(prs, sw, sh):
    s = add_blank_slide(prs)
    add_header(s, sw, "RESULTS  /  MODEL COMPARISON",
               "Model Performance Visualized")
    chart_path = render_model_comparison_chart()
    add_image_fit(s, chart_path, Inches(0.5), Inches(1.7),
                  sw - Inches(1.0), Inches(4.7))

    # Takeaway caption under chart (simple text, no card)
    cap = s.shapes.add_textbox(Inches(0.5), Inches(6.05),
                               sw - Inches(1.0), Inches(0.4))
    tf = cap.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("TAKEAWAY  ·  Random Forest baseline dominates on R² and MAE; "
              "OLS fails on MAPE — the linear comparator confirms non-linearity.")
    p.font.name = "Inter"; p.font.size = Pt(11); p.font.color.rgb = NAVY
    p.font.bold = True
    add_speaker_strip(s, sw,
        "The visual story: OLS is on a different scale of error entirely. RF and "
        "XGBoost are close on MAPE, but RF wins on R² and MAE — that's why we deployed it.")
    set_notes(s,
        "Three-panel chart shows R-squared (left), MAPE (middle, OLS clipped), MAE (right). "
        "Random Forest in amber is the deployed model. OLS in gray is the linear baseline. "
        "XGBoost in navy is the tuned alternative. Key talking points: RF baseline beat tuned "
        "RF and tuned XGBoost — Decision 25 confirmation grid verified this. Sample size "
        "(1,212 train rows) is the binding constraint on tuning gains, not configuration.")
    return s


def build_rf_scatter_slide(prs, sw, sh):
    s = add_blank_slide(prs)
    add_header(s, sw, "RESULTS  /  PREDICTION QUALITY",
               "Random Forest: Actual vs. Predicted")

    img = EDA / "rf_actual_vs_predicted.png"
    add_image_fit(s, img, Inches(0.5), Inches(1.7),
                  Inches(8.3), Inches(4.7))

    txt_left = Inches(9.0)
    tb = s.shapes.add_textbox(txt_left, Inches(1.8),
                              Inches(4.0), Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True

    items = [
        ("Tightest cluster < PHP 10M",
         "where listing data is densest"),
        ("Wider dispersion > PHP 20M",
         "luxury tier — listings thin out"),
        ("299 held-out test rows",
         "random_state = 42"),
        ("Diagonal reference",
         "perfect-prediction line"),
    ]
    for i, (head, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = head
        p.font.name = "Inter"; p.font.size = Pt(13)
        p.font.bold = True; p.font.color.rgb = NAVY
        p.space_before = Pt(6) if i else Pt(0)

        p = tf.add_paragraph(); p.text = body
        p.font.name = "Inter"; p.font.size = Pt(11); p.font.color.rgb = GRAY
        p.space_after = Pt(2)

    add_speaker_strip(s, sw,
        "The cluster shape is the story: dense and on-diagonal where data is thick, "
        "wider where the luxury tier dilutes the sample. This is the per-row picture "
        "behind the 0.807 R-squared.")
    set_notes(s,
        "RF actual vs predicted scatter on 299 held-out test properties. The diagonal is "
        "perfect prediction. Note the dense, on-diagonal cluster below PHP 10M — this is "
        "where the model is most reliable. Above PHP 20M the dispersion widens; that's the "
        "luxury tier where listings thin out and asking-price noise dominates. This visualizes "
        "what the per-LGU MAPE table tells in numbers.")
    return s


def build_per_lgu_chart_slide(prs, sw, sh):
    s = add_blank_slide(prs)
    add_header(s, sw, "RESULTS  /  SEGMENT EVALUATION",
               "Per-LGU MAPE Visualized")

    chart_path = render_per_lgu_mape_chart()
    add_image_fit(s, chart_path, Inches(0.5), Inches(1.7),
                  Inches(8.5), Inches(4.7))

    # Right column callouts
    add_card(s, Inches(9.2), Inches(1.8), Inches(3.8), Inches(1.4),
             "MOST CREDIBLE",
             "Lapu-Lapu  ·  40.0%",
             "n=108. Largest sample, best-evidenced LGU benchmark.")

    add_card(s, Inches(9.2), Inches(3.4), Inches(3.8), Inches(1.4),
             "CAUTION  ·  SMALL SAMPLE",
             "Talisay (n=13) · Minglanilla (n=12)",
             "Reported as indicative; not reliable single-LGU estimates.",
             accent=RGBColor(0xB5, 0xBC, 0xC8))

    add_card(s, Inches(9.2), Inches(5.0), Inches(3.8), Inches(1.4),
             "PATTERN",
             "Sample size drives credibility",
             "Smaller LGUs sit at both extremes — error follows n, not geography.")

    add_speaker_strip(s, sw,
        "Visualized, the per-LGU pattern is clear: error follows sample size before it "
        "follows location. Lapu-Lapu at 108 properties is our anchor; Talisay and "
        "Minglanilla are flagged as indicative.")
    set_notes(s,
        "Horizontal bar chart sorted ascending. Amber bar is Lapu-Lapu — most credible "
        "benchmark with n=108. Gray bars are small samples (Talisay n=13, Minglanilla n=12) "
        "where confidence intervals would be wide. Navy bars are adequate samples. "
        "Source: chapter7_eval_summary_2026-05-05.json — per-segment evaluation supplement.")
    return s


def build_rq_answers_slide(prs, sw, sh):
    s = add_blank_slide(prs)
    add_header(s, sw, "CONCLUSION  /  RESEARCH QUESTIONS",
               "Research Questions Answered")

    # 2x2 grid of cards
    card_w = Inches(6.2)
    card_h = Inches(2.3)
    gap_x = Inches(0.3)
    gap_y = Inches(0.25)
    left0 = Inches(0.5)
    top0 = Inches(1.7)

    cards = [
        ("RQ 1  ·  ANSWERED",
         "Value drivers identified",
         "Structural features dominate (75.6% SHAP) — property type, area, "
         "bedrooms. CBD distance block (21.4%) is the largest geospatial signal, "
         "led by dist_consolacion. MCRAI recreation enters the composite."),
        ("RQ 2  ·  ANSWERED",
         "Random Forest yields lowest error",
         "RF baseline: R²=0.807, MAPE=59.28%, MAE=PHP 4.95M. "
         "Beats OLS (R²=0.083) by +0.72 and tuned XGBoost (R²=0.557). "
         "Decision 25 confirmation: tuning does not improve at this sample size."),
        ("RQ 3  ·  ANSWERED",
         "Geospatial features improve performance",
         "CBD distance block alone contributes 21.4% of total SHAP weight — "
         "more than every MCRAI category combined. Network distance via osmnx "
         "is essential for Mactan island properties (n=312)."),
        ("RQ 4  ·  ANSWERED",
         "Valuation gap quantified",
         "BIR zonal systematically below RF prediction in fast-moving "
         "barangays. Gap magnitude flagged as evidence base for LGU zonal "
         "revision — directly actionable output."),
    ]
    for i, (eyebrow, title, body) in enumerate(cards):
        col = i % 2; row = i // 2
        left = left0 + col * (card_w + gap_x)
        top = top0 + row * (card_h + gap_y)
        add_card(s, left, top, card_w, card_h, eyebrow, title, body)

    add_speaker_strip(s, sw,
        "All four research questions answered with direct evidence: feature importance "
        "for RQ1, model comparison for RQ2, geospatial contribution for RQ3, "
        "and the BIR-versus-model gap for RQ4.")
    set_notes(s,
        "Walk through each card briefly. RQ1 is answered by SHAP analysis. RQ2 by the "
        "model comparison table. RQ3 by the SHAP block analysis showing 21.4% CBD distance "
        "contribution. RQ4 by the valuation gap map. If the panel asks 'did you answer your "
        "questions?' — this slide is the direct response.")
    return s


def build_synthesis_slide(prs, sw, sh):
    s = add_blank_slide(prs)
    add_header(s, sw, "CONCLUSION  /  SYNTHESIS",
               "What This Thesis Proves")

    qb = s.shapes.add_textbox(Inches(1.0), Inches(2.0),
                              sw - Inches(2.0), Inches(3.8))
    tf = qb.text_frame; tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = ("This thesis proves that property-level, geospatially grounded "
              "valuation is feasible in Metro Cebu using only publicly "
              "available data.")
    p.font.name = "Inter"; p.font.size = Pt(22)
    p.font.color.rgb = NAVY; p.font.bold = True
    p.space_after = Pt(18)

    p = tf.add_paragraph()
    p.text = "The 59% MAPE is the honest ceiling of asking-price data."
    p.font.name = "Inter"; p.font.size = Pt(18); p.font.color.rgb = DARK
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = ("The per-sqm MAE of one BIR zonal tier is the honest floor "
              "of practical utility.")
    p.font.name = "Inter"; p.font.size = Pt(18); p.font.color.rgb = DARK
    p.space_after = Pt(18)

    p = tf.add_paragraph()
    p.text = "Defensible, reproducible, and ready for practitioners today."
    p.font.name = "Inter"; p.font.size = Pt(18); p.font.italic = True
    p.font.color.rgb = AMBER; p.font.bold = True

    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(1.0), Inches(6.4),
                                Inches(0.9), Inches(0.06))
    accent.line.fill.background()
    accent.fill.solid(); accent.fill.fore_color.rgb = AMBER

    footer = s.shapes.add_textbox(Inches(1.0), Inches(6.55),
                                  sw - Inches(2.0), Inches(0.4))
    p = footer.text_frame.paragraphs[0]
    p.text = "Synthesis text  ·  copy-paste into the closing slide as the lead paragraph."
    p.font.name = "Inter"; p.font.size = Pt(10); p.font.italic = True
    p.font.color.rgb = GRAY
    return s


def build_recommendations_slide(prs, sw, sh):
    s = add_blank_slide(prs)
    add_header(s, sw, "CONCLUSION  /  RECOMMENDATIONS",
               "Actionable Recommendations")

    # 5 cards in 2x3 grid (last cell wide for future research)
    card_w = Inches(4.05)
    card_h = Inches(2.15)
    gap_x = Inches(0.2)
    gap_y = Inches(0.25)
    left0 = Inches(0.5)
    top0 = Inches(1.7)

    top_cards = [
        ("FOR LGUs",
         "Use the model–BIR gap as zonal-revision evidence",
         "Prioritize barangays where RF prediction exceeds BIR zonal by >30%. "
         "Real property tax base is undervalued where the gap is widest."),
        ("FOR BANKS",
         "Use as a collateral pre-screen",
         "Flag listings or appraisals diverging by more than one BIR zonal "
         "tier (~PHP 19,743/sqm) for human review."),
        ("FOR BUYERS / SELLERS",
         "Reference price with feature-level explanation",
         "Use the SHAP waterfall to understand which features drive the "
         "estimate before negotiating or pricing."),
    ]
    for i, (eyebrow, title, body) in enumerate(top_cards):
        left = left0 + i * (card_w + gap_x)
        add_card(s, left, top0, card_w, card_h, eyebrow, title, body)

    bottom_cards = [
        ("FOR URBAN PLANNERS",
         "Map accessibility premiums spatially",
         "Use the predicted price surface to target where infrastructure "
         "(CBRT, MCX, SRP) generates the largest residential premium."),
        ("FOR FUTURE RESEARCHERS",
         "Integrate transaction-grade data",
         "Registry of Deeds + BIR eFPS records remove the asking-price noise "
         "floor. Build a temporal panel; develop a separate vacant-lot model."),
    ]
    bottom_w = Inches(6.2)
    for i, (eyebrow, title, body) in enumerate(bottom_cards):
        left = left0 + i * (bottom_w + gap_x)
        add_card(s, left, top0 + card_h + gap_y, bottom_w, card_h,
                 eyebrow, title, body)

    add_speaker_strip(s, sw,
        "Five concrete next steps targeted to the people who would actually use this — "
        "LGUs, banks, buyers, planners, and the next round of researchers. "
        "Decision-support, with clear thresholds and integration paths.")
    set_notes(s,
        "These are declarative recommendations, not just use cases. Each one names a "
        "specific stakeholder and a specific action. Note the LGU recommendation has a "
        "concrete threshold (>30%) and the bank recommendation references the per-sqm "
        "MAE benchmark. Future research recommendations name specific data sources.")
    return s


# ----------------------- reorder helper -----------------------
def move_slide(prs, from_idx, to_idx):
    """Move slide element from from_idx to to_idx in <sldIdLst>."""
    sldIdLst = prs.slides._sldIdLst
    elements = list(sldIdLst)
    elem = elements[from_idx]
    sldIdLst.remove(elem)
    sldIdLst.insert(to_idx, elem)


# ----------------------- main -----------------------
def main():
    prs = Presentation(str(SRC))
    sw = prs.slide_width
    sh = prs.slide_height
    n_orig = len(prs.slides)
    print(f"Original slide count: {n_orig}")
    print(f"Size: {sw/914400:.2f} x {sh/914400:.2f} in")

    # Build new slides — they get appended at the end.
    # Track each new slide's index after appending.
    s_perf = build_model_perf_slide(prs, sw, sh)        # idx n_orig+0
    s_scatter = build_rf_scatter_slide(prs, sw, sh)     # idx n_orig+1
    s_per_lgu = build_per_lgu_chart_slide(prs, sw, sh)  # idx n_orig+2
    s_rq = build_rq_answers_slide(prs, sw, sh)          # idx n_orig+3
    s_rec = build_recommendations_slide(prs, sw, sh)    # idx n_orig+4
    s_syn = build_synthesis_slide(prs, sw, sh)          # idx n_orig+5 (stays at end)

    # Reorder. Move new slides one at a time to their target positions.
    # Each move is from current append-end to a target index.
    # IMPORTANT: do moves from EARLIEST target → LATEST so later targets
    # remain valid (they refer to positions in the still-evolving list).
    # 0-indexed targets:
    #   Model Perf  → idx 26  (becomes position 27)
    #   RF Scatter  → idx 28  (becomes position 29)
    #   Per-LGU     → idx 33  (becomes position 34)
    #   RQ Answers  → idx 40  (becomes position 41)
    #   Recommend.  → idx 42  (becomes position 43)

    # After each move, the appended slides shift. We track them by name.
    # Because move_slide pops from from_idx then inserts at to_idx, doing
    # earliest-first means later appended slides drift down by one index
    # each time. Account for that.

    # Initial append indices (0-based) for the 5 new slides:
    #   perf = n_orig+0 = 39
    #   scat = n_orig+1 = 40
    #   plgu = n_orig+2 = 41
    #   rq   = n_orig+3 = 42
    #   rec  = n_orig+4 = 43

    # Step 1: move perf (idx 39) → 26.
    # All later appended slides (40,41,42,43) stay at same indices (move
    # pulled an earlier index up; indices > 26 each shift +1 after insert,
    # but the source index 39 also shifted to 40 after... actually let me
    # think carefully: removing at 39 then inserting at 26 — elements at
    # idx 26..38 each move +1; element at 39 gets re-inserted at 26;
    # elements at 40..43 stay at 40..43.
    move_slide(prs, n_orig + 0, 26)
    # Now: scat is still at 40, plgu at 41, rq at 42, rec at 43.

    # Step 2: move scat (idx 40) → 28.
    # After: elements 28..39 shift +1; plgu was at 41 → 41, rq at 42 → 42,
    # rec at 43 → 43. (only items between 28 and source-1 shift.)
    move_slide(prs, n_orig + 1, 28)
    # Now: plgu at 41, rq at 42, rec at 43.

    # Step 3: move plgu (idx 41) → 33.
    move_slide(prs, n_orig + 2, 33)
    # Now: rq at 42, rec at 43.

    # Step 4: move rq (idx 42) → 40.
    move_slide(prs, n_orig + 3, 40)
    # Now: rec at 43.

    # Step 5: move rec (idx 43) → 42.
    move_slide(prs, n_orig + 4, 42)

    prs.save(str(SRC))
    print(f"Final slide count: {len(prs.slides)}")
    print(f"Saved in place: {SRC.name}")


if __name__ == "__main__":
    main()
