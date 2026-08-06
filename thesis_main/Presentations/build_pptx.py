"""Build defense_deck.pptx — mirrors defense_deck.html (same content + deep-green design).
Run:  python3 build_pptx.py
Photo slides use gray placeholder frames; drop real images in PowerPoint.
two_beat second line + claim reveals: add click animations yourself in PowerPoint.
"""
import os, re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
DG, EDA, DECK = "../Manuscript/diagrams/", "../EDA/plots/", "assets/deck/"
def R(p): return os.path.join(HERE, p)

# palette
BG=RGBColor(0xFA,0xF8,0xF4); INK=RGBColor(0x1A,0x1A,0x1A); MUTED=RGBColor(0x6B,0x6B,0x66)
ACCENT=RGBColor(0x2F,0x5D,0x50); SOFT=RGBColor(0xE7,0xEE,0xEA); LINE=RGBColor(0xD9,0xD6,0xCF)
WHITE=RGBColor(0xFF,0xFF,0xFF); GREY=RGBColor(0x9A,0x96,0x8C); SUPP=RGBColor(0x33,0x33,0x2F)
FONT="Helvetica Neue"
EW, EH = 13.333, 7.5

prs = Presentation(); prs.slide_width=Inches(EW); prs.slide_height=Inches(EH)
BLANK = prs.slide_layouts[6]

def entity(s):
    return (s.replace("&amp;","&").replace("&rarr;","→").replace("&Sigma;","Σ")
             .replace("&isin;","∈").replace("&beta;","β").replace("&nbsp;"," ")
             .replace("<sub>","").replace("</sub>","").replace("<sup>","").replace("</sup>",""))
def parse_runs(text):
    out=[]
    for p in re.split(r"(<strong>.*?</strong>)", text):
        if not p: continue
        if p.startswith("<strong>"): out.append((entity(p[8:-9]), True))
        else: out.append((entity(p), False))
    return out

def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s

def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    return tf
def para(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT, first=False, ls=1.06, after=4, bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; p.line_spacing=ls; p.space_after=Pt(after)
    if bullet:
        r=p.add_run(); r.text="▪  "; r.font.size=Pt(size); r.font.color.rgb=ACCENT; r.font.name=FONT; r.font.bold=True
    for txt, b in (parse_runs(text) if isinstance(text,str) else [(text,bold)]):
        r=p.add_run(); r.text=txt; f=r.font; f.size=Pt(size); f.color.rgb=color; f.bold=(b or bold); f.name=FONT
    return p
def kicker(s, text, l=0.9, t=0.62, color=ACCENT, align=PP_ALIGN.LEFT, w=11.5):
    para(box(s,l,t,w,0.6), text.upper(), 15, color, True, align, first=True)
def cite(s, c):
    if c: para(box(s,6.0,EH-0.62,6.43,0.45,), c, 12, MUTED, align=PP_ALIGN.RIGHT, first=True)
def img_fit(s, path, bx, by, bw, bh):
    iw,ih = Image.open(R(path)).size; sc=min(bw/iw, bh/ih); w,h=iw*sc, ih*sc
    s.shapes.add_picture(R(path), Inches(bx+(bw-w)/2), Inches(by+(bh-h)/2), Inches(w), Inches(h))
def rect(s, l, t, w, h, fill, line=None, dash=False):
    sp=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False
    return sp

# ---- slide builders -------------------------------------------------------
def title():
    s=slide()
    para(box(s,0.9,2.05,11.5,0.6), "UA&P · BS Data Science · June 2026", 16, ACCENT, True, PP_ALIGN.CENTER, first=True)
    tf=box(s,0.9,2.5,11.5,2.4,MSO_ANCHOR.TOP)
    para(tf,"Predicting Open-Market Residential Property Values in Metro Cebu",40,INK,True,PP_ALIGN.CENTER,first=True,ls=1.08)
    para(box(s,0.9,4.95,11.5,0.6),"Using Machine Learning and Geospatial Features",22,MUTED,align=PP_ALIGN.CENTER,first=True)
    tf2=box(s,0.9,5.7,11.5,1.0)
    para(tf2,"Chris Dominic Estreba",18,MUTED,align=PP_ALIGN.CENTER,first=True)
    para(tf2,"BS Data Science Capstone",18,MUTED,align=PP_ALIGN.CENTER)

def divider(num, name, sub=""):
    s=slide(ACCENT)
    tf=box(s,0.9,2.6,11.5,2.2,MSO_ANCHOR.MIDDLE)
    para(tf,num,22,SOFT,True,PP_ALIGN.CENTER,first=True,after=10)
    para(tf,name,56,WHITE,True,PP_ALIGN.CENTER)
    if sub: para(tf,sub,18,SOFT,align=PP_ALIGN.CENTER)

def picture(head, sub="", q="", c="", photo=""):
    s=slide()
    fr=rect(s,0.85,1.1,6.1,5.3,SOFT);
    tf=fr.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    para(tf,f"[ photo: {photo} ]",16,GREY,align=PP_ALIGN.CENTER,first=True)
    body=box(s,7.3,1.1,5.3,5.3,MSO_ANCHOR.MIDDLE)
    para(body,head,34,INK,True,first=True,ls=1.12)
    if sub: para(body,sub,22,MUTED,after=2)
    if q: para(body,q,28,ACCENT,True)
    cite(s,c)

def statement(text, sub="", kicker_t="", c="", stat="", align="center"):
    s=slide(); al = PP_ALIGN.CENTER if align=="center" else PP_ALIGN.LEFT
    tf=box(s,0.9,1.4,11.5,4.7,MSO_ANCHOR.MIDDLE)
    if kicker_t: para(tf,kicker_t.upper(),15,ACCENT,True,al,first=True,after=14)
    if stat: para(tf,stat,80,ACCENT,True,al,first=(not kicker_t),after=6)
    para(tf,text,34,INK,True,al,first=(not kicker_t and not stat),ls=1.14)
    if sub: para(tf,sub,23,MUTED,al,after=2)
    cite(s,c)

def claim_support(claim, bullets, kicker_t="", c="", stat="", sub=""):
    s=slide(); top=1.25
    tf=box(s,0.9,top,11.5,2.2)
    fst=True
    if kicker_t: para(tf,kicker_t.upper(),15,ACCENT,True,first=True,after=10); fst=False
    if stat: para(tf,stat,72,ACCENT,True,first=fst,after=2); fst=False
    para(tf,claim,38,INK,True,first=fst,ls=1.12)
    if sub: para(tf,sub,22,MUTED)
    bt=box(s, 0.9, 3.9, 11.5, 3.2)
    for i,b in enumerate(bullets): para(bt,b,24,SUPP,first=(i==0),bullet=True,after=10,ls=1.3)
    cite(s,c)

def points(kicker_t, items, c=""):
    s=slide(); kicker(s,kicker_t)
    bt=box(s,0.9,1.7,11.5,5.0,MSO_ANCHOR.MIDDLE)
    for i,b in enumerate(items): para(bt,b,27,INK,first=(i==0),bullet=True,after=14,ls=1.3)
    cite(s,c)

def plot(kicker_t, img, caption="", c="", contain=True):
    s=slide(); kicker(s,kicker_t)
    img_fit(s,img,0.9,1.55,11.5,4.5)
    if caption: para(box(s,0.9,6.2,11.5,0.7),caption,21,INK,align=PP_ALIGN.CENTER,first=True)
    cite(s,c)

def insight(kicker_t, img, bullets, c=""):
    s=slide(); kicker(s,kicker_t)
    img_fit(s,img,0.7,1.6,6.6,5.0)
    bt=box(s,7.6,1.6,5.1,5.0,MSO_ANCHOR.MIDDLE)
    for i,b in enumerate(bullets): para(bt,b,23,SUPP,first=(i==0),bullet=True,after=16,ls=1.32)
    cite(s,c)

def split(kicker_t, img1, img2, cap1="", cap2="", c=""):
    s=slide(); kicker(s,kicker_t)
    img_fit(s,img1,0.7,1.6,5.7,4.3); img_fit(s,img2,6.9,1.6,5.7,4.3)
    if cap1: para(box(s,0.7,6.0,5.7,0.8),cap1,18,INK,align=PP_ALIGN.CENTER,first=True)
    if cap2: para(box(s,6.9,6.0,5.7,0.8),cap2,18,INK,align=PP_ALIGN.CENTER,first=True)
    cite(s,c)

def formula(kicker_t, name, tex, sub, c=""):
    s=slide()
    tf=box(s,0.9,1.6,11.5,4.4,MSO_ANCHOR.MIDDLE)
    para(tf,kicker_t.upper(),15,ACCENT,True,PP_ALIGN.CENTER,first=True,after=16)
    para(tf,name,32,INK,True,PP_ALIGN.CENTER,after=18,ls=1.1)
    fr=rect(s,2.4,3.45,8.5,1.0,SOFT); ftf=fr.text_frame; ftf.word_wrap=True; ftf.vertical_anchor=MSO_ANCHOR.MIDDLE
    para(ftf,entity(tex),26,INK,align=PP_ALIGN.CENTER,first=True)
    para(box(s,0.9,4.7,11.5,0.8),sub,22,MUTED,align=PP_ALIGN.CENTER,first=True)
    cite(s,c)

def flow(kicker_t, stages, c=""):
    s=slide(); kicker(s,kicker_t)
    tf=box(s,0.9,2.9,11.5,1.6,MSO_ANCHOR.MIDDLE)
    para(tf,"   →   ".join(stages),26,ACCENT,True,PP_ALIGN.CENTER,first=True,ls=1.3)
    cite(s,c)

def qa(question, answer, sub="", c=""):
    s=slide()
    ql=box(s,0.9,1.6,5.0,4.3,MSO_ANCHOR.MIDDLE)
    para(ql,question,36,ACCENT,True,first=True,ls=1.15)
    rect(s,6.05,2.2,0.03,3.1,ACCENT)
    al=box(s,6.4,1.6,6.0,4.3,MSO_ANCHOR.MIDDLE)
    para(al,answer,40,INK,True,first=True,ls=1.14)
    if sub: para(al,sub,22,MUTED,after=2)
    cite(s,c)

def two_beat(first, second, c=""):
    s=slide()
    tf=box(s,0.9,1.6,11.5,4.3,MSO_ANCHOR.MIDDLE)
    para(tf,first,40,INK,True,PP_ALIGN.CENTER,first=True,after=34,ls=1.14)
    para(tf,second,46,ACCENT,True,PP_ALIGN.CENTER,ls=1.12)
    para(box(s,0.9,6.55,11.5,0.5),"▶ animate the second line on click",13,MUTED,align=PP_ALIGN.CENTER,first=True)
    cite(s,c)

def add_table(s, headers, rows, top, ratios):
    ncols=len(headers); nrows=len(rows)+1; width=11.5
    gt=s.shapes.add_table(nrows,ncols,Inches(0.9),Inches(top),Inches(width),Inches(0.5*nrows)).table
    gt.first_row=False; gt.horz_banding=False
    tot=sum(ratios)
    for j,rt in enumerate(ratios): gt.columns[j].width=Inches(width*rt/tot)
    def setcell(cell,text,color,fill,size,bold,align):
        cell.fill.solid(); cell.fill.fore_color.rgb=fill
        cell.margin_left=Inches(0.12); cell.margin_top=Inches(0.05); cell.margin_bottom=Inches(0.05)
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        tf=cell.text_frame; p=tf.paragraphs[0]; p.alignment=align
        r=p.add_run(); r.text=str(text); r.font.size=Pt(size); r.font.color.rgb=color; r.font.bold=bold; r.font.name=FONT
    for j,h in enumerate(headers):
        setcell(gt.cell(0,j),h,WHITE,ACCENT,15,True,PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)
    for i,row in enumerate(rows,1):
        fill = SOFT if i%2==0 else WHITE
        for j,v in enumerate(row):
            setcell(gt.cell(i,j),v,INK if j==0 else ACCENT,fill,18 if j==0 else 19, j==0, PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)

def table(kicker_t, headers, rows, c=""):
    s=slide(); kicker(s,kicker_t); add_table(s,headers,rows,2.1,[1.4,1,1]); cite(s,c)

def feature_matrix(kicker_t, headers, rows, c=""):
    s=slide(); kicker(s,kicker_t); add_table(s,headers,rows,1.55,[2.3,1,1,1]); cite(s,c)

def close(lead, accent_line):
    s=slide()
    tf=box(s,0.9,1.6,11.5,4.3,MSO_ANCHOR.MIDDLE)
    para(tf,lead,30,INK,align=PP_ALIGN.CENTER,first=True,ls=1.25,after=16)
    para(tf,accent_line,40,ACCENT,True,PP_ALIGN.CENTER,ls=1.14)

def references(items):
    s=slide(); kicker(s,"References")
    half=(len(items)+1)//2
    for k,(l,group) in enumerate([(0.9,items[:half]),(6.9,items[half:])]):
        tf=box(s,l,1.4,5.7,5.7)
        for i,it in enumerate(group): para(tf,it,11,SUPP,first=(i==0),after=7,ls=1.2)

def thankyou():
    s=slide(ACCENT)
    tf=box(s,0.9,2.7,11.5,2.0,MSO_ANCHOR.MIDDLE)
    para(tf,"Thank you",58,WHITE,True,PP_ALIGN.CENTER,first=True,after=12)
    para(tf,"Chris Dominic Estreba · Questions welcome",20,SOFT,align=PP_ALIGN.CENTER)

# ---- the deck (mirrors build_deck.py) -------------------------------------
title()
picture("Everyone needs a place to live.", photo="a Cebu family / a home")
picture("A home is the biggest purchase most families ever make.", photo="Cebu homes / skyline")
picture("You find one. ₱6.5M.", q="Is that fair?", photo="a property listing")
picture("Ask a broker — you get a number.", photo="a real-estate broker")
picture("Ask an appraiser — you get a different one.", photo="an appraiser at work")
statement("None of them agree. And none of them says why.")
picture("Now you're the seller. What is it worth?", photo="a 'For Sale' sign")
picture("Too high, it sits for months. Too low, you leave money behind.", photo="an empty house / keys")
picture("Everyone guesses — from a different corner.", photo="one house, five question marks (buyer · seller · broker · bank · LGU)")
picture("BIR zonal values — built for taxation, not live market prices. And not always updated.", c="BIR, n.d.; Otsuka et al., 2023", photo="a tax / zonal schedule")
picture("Bank appraisals — tied to lending and collateral risk. A different question than market value.", photo="a bank facade")
picture("Online listings — asking prices, not verified sales. Full of seller strategy and noise.", photo="a property listings website")
two_beat("Each one is useful. But they answer different questions — and they don't line up.","The references exist. Nothing connects them.")
divider("02","The Market")
claim_support("Cebu's economy grew in 2024.",["Among the fastest-growing provinces in the country.","Driven by IT-BPM and call-center expansion.","And a strong post-pandemic rebound in tourism."],stat="7.3%",c="PSA Cebu, 2025")
claim_support("And home prices followed in 2025.",["Metro Cebu's growth — the highest rate outside Metro Manila.","From the BSP residential property price index (RREPI).","Prices moving faster than benchmarks and incomes can adjust."],stat="11.5%",c="BSP, 2025")
picture("New infrastructure keeps redrawing what counts as ‘near’.", photo="CBRT / Cebu–Cordova Expressway / SRP")
divider("03","The Question")
statement("No defensible, market-facing price reference for Metro Cebu homes.",kicker_t="The business problem")
claim_support("Turn a property — and where it sits — into a defensible price per square meter.",["Property-level data at scale, drawn from online listings.","Location turned into measurable features — distances, accessibility, neighborhood price.","A model that generalizes honestly, and explains its reasoning."],kicker_t="The data problem")
points("Objectives",["<strong>Predictive</strong> — estimate open-market price per square meter across Metro Cebu.","<strong>Prescriptive</strong> — show where those estimates diverge from official benchmarks."])
points("Research Questions",["<strong>1</strong>   What value drivers influence price?","<strong>2</strong>   Which model deploys best?","<strong>3</strong>   Do geospatial features help?","<strong>4</strong>   How large is the BIR gap?"])
statement("A stratified, explainable ML model — delivered as a prototype for triangulation.",kicker_t="The solution")
flow("Roadmap — CRISP-DM",["Business","Data","Prep","Modeling","Evaluation","Deployment"])
divider("04","The Field")
points("Two traditions",["<strong>Hedonic regression</strong> — interpretable, classic.","<strong>Machine learning</strong> — flexible, higher accuracy."],c="Rosen, 1974; Breiman, 2001")
points("What drives price in the literature",["Structural","Economic / benchmark","Geospatial / accessibility","Amenity / point-of-interest"])
points("International evidence",["Tree models tend to outperform linear models on tabular property data.","But small samples demand caution."],c="Grinsztajn et al., 2022; Tanamal et al., 2023")
points("Philippine & Cebu evidence",["ML work is mostly Manila / Pangasinan-centric.","<strong>Agosto (2020)</strong> is the only Cebu study — transport accessibility is the primary driver."],c="Viray, 2023; Ramolete et al., 2023; Agosto, 2020")
statement("Nothing property-level, geospatial, and explainable — for Metro Cebu.",kicker_t="The gap")
claim_support("This study builds exactly that.",["Property-level — every listing geocoded across the six LGUs.","Open-market — aligned to the IVS Market Value basis.","Geospatial and explainable — with SHAP behind every prediction."],kicker_t="Our bridge")
qa("Why open-market listings?","The closest available proxy to IVS Market Value.",sub="Arm's-length asking evidence — the basis a willing buyer and seller actually meet on.",c="IVSC, 2025")
qa("Why geospatial features?","Location is price — so we construct the location signal, not ingest it.")
divider("05","The Data")
insight("Collection funnel",DECK+"funnel_collection.png",["16,561 listings scraped — only <strong>3,616</strong> survived.","Geocoding, the six-LGU and residential filters, a price-per-sqm sanity band, and de-duplication did the cutting.","OnePropertee was dropped entirely — city-centroid geocoding and mis-extracted prices."])
plot("Study area",DG+"lgu_boundaries.png","The six LGUs in scope.")
split("Where the listings sit",DG+"properties_by_stratum.png",DECK+"listings_by_lgu.png",cap1="Across the six LGUs.",cap2="Concentrated in Cebu City and the island; thin in the south.")
plot("One row of the data",DECK+"abt_snapshot_wide.png","18 of 51 columns — structural, location, distances, MCRAI, benchmarks.")
plot("Engineered location — CBD distance",DG+"study_area_clean.png","Shortest-path road distance to 8 polycentric nodes.",c="Giuliano & Small, 1991; McMillen, 2003; JICA, 2015; Boeing, 2017")
formula("MCRAI — what it is","Metro Cebu Residential Accessibility Index (MCRAI)","MCRAIic = Σ 1 / max(dij, 0.5)²   (β = 2)","Nearer amenities count more — accessibility decays with the square of road-distance.",c="Hansen, 1959")
plot("MCRAI — how it's built",DG+"amenities_map.png","Eight amenity categories, each with its own radius. The composite keeps only positive-signal categories.")
insight("MCRAI — how one property scores",DECK+"mcrai_catchment.png",["Each category has its own catchment radius (1–5 km).","Nearer amenities count more — inverse-square decay (β = 2).","Transport enters separately, as road-network distance."],c="Hansen, 1959")
plot("Engineered location — spatial lag",EDA+"09_data_integrity/Master_geocoding_clusters.png","Mean price of nearby same-type listings within 500 m.",c="Tobler, 1970")
insight("Price per square meter is skewed",EDA+"01_target/all_strata_price_boxplot.png",["Price per square meter is right-skewed — a few premium units stretch the tail.","Modeled in <strong>log</strong>: it tames the skew but keeps a per-unit meaning.","Predictions are back-transformed for the price surface."])
insight("Price varies across the map",EDA+"02_geographic/price_by_lgu_faceted.png",["Premium signal sits in the core and island: Cebu City ~₱113,600/sqm, Mandaue ~₱96,100, Lapu-Lapu ~₱92,100.","Talisay, Minglanilla, Consolacion form a lower band (~₱47–57k).","Location can't be assumed uniform — it has to be modeled."])
insight("Do we trust the features?",EDA+"04_correlation/Houses_feature_correlation_heatmap.png",["Several CBD-distance features move together — the southern corridor is geographically nested.","A flag for the linear baseline, not the trees: OLS trimmed the collinear terms (confirmed with VIF).","Random Forest tolerates correlated predictors, so the meaningful nodes stayed."])
divider("06","Groundwork")
insight("The finding that changed everything",EDA+"01_target/all_strata_price_boxplot.png",["The condo median runs ≈ <strong>5.8×</strong> the vacant-lot median.","Built area vs land-only are different price logics.","A single pooled model would blur them — so we stratify by type."])
claim_support("Cleaning kept every decision visible.",["Imputed values were flagged, never silently filled.","Structurally-absent fields (beds/baths for lots) left blank, not faked.","Hard duplicates dropped; a price-per-sqm sanity band applied."],kicker_t="Cleaning, honestly")
claim_support("One feature was quietly hijacking the target.",["A scale-selector feature dominated the SHAP rankings.","The target had come to mean two different things across scrape batches.","Tracing and fixing it is why the final metrics can be trusted."],kicker_t="A data-integrity story")
divider("07","The Build")
statement("A single model would treat condos, houses, and lots as one market — and they differ 5.8×.",kicker_t="One model isn't enough")
points("So we built three",["<strong>Condominium</strong> — 1,300","<strong>Houses</strong> — 1,223","<strong>Vacant Lot</strong> — 849"],c="Dröes et al., 2019; Usman et al., 2020")
feature_matrix("Features per stratum",["Feature group","Condo","Houses","Lot"],[("Structural — area","✓","✓","✓"),("Bedrooms / bathrooms","✓","✓","—"),("CBD distances (8 nodes)","✓","✓","✓"),("MCRAI composite","✓","✓","—"),("MCRAI individual categories","—","—","✓ (6)"),("BIR zonal + spatial lag","✓","✓","✓"),("Property-type dummies","—","✓","—"),("Total features","21","24","22")])
points("The model lineup",["<strong>OLS</strong> — interpretable hedonic baseline.","<strong>Random Forest & XGBoost</strong> — capture non-linear, interacting effects.","<strong>SHAP</strong> — explains every prediction, feature by feature."],c="Rosen, 1974; Breiman, 2001; Chen & Guestrin, 2016; Lundberg & Lee, 2017")
claim_support("GroupKFold by location — the honest test.",["Listings are grouped by coordinate cluster.","The same spot never appears in both train and test.","Stricter than a random 80/20 split — it blocks neighborhood memorization."],kicker_t="Honest evaluation")
points("Random Forest deployed (RQ2)",["Tree models beat OLS in every stratum.","RF edged XGBoost (19.3 vs 19.8 / 22.7 vs 23.6 / 38.4 vs 40.2 MdAPE).","Deployed for robustness on small samples and simplicity."])
divider("08","The Appraisal")
table("Headline accuracy",["Property type","Typical error (MdAPE)","Within 20% (PE20)"],[["Condominium","19.3%","51%"],["Houses","22.7%","44%"],["Vacant Lot","38.4%","26%"]])
points("What that means",["<strong>MdAPE</strong> — half of estimates fall within this error.","<strong>PE20</strong> — the share landing within 20% (the practical hit-rate).","Reported with MAPE, COD, PRD — but not claiming IAAO compliance."])
insight("Do geospatial features help? (RQ3)",DECK+"ablation_tiers.png",["Geospatial features improve <strong>every</strong> stratum over structural-only.","Biggest lift where benchmarks are weakest: vacant lots (+13 pts) and condos (+5.7).","For houses, city + BIR zonal already capture most of the location signal."])
insight("What drives price — condominiums (RQ1)",EDA+"10_stratified_models/shap_condo_rf_bar.png",["Neighborhood price level (spatial lag) leads — condos track their block.","Weight spreads across several nodes — Consolacion, airport, Mactan, CBP.","The signature of a multi-node, polycentric market."])
insight("What drives price — houses (RQ1)",EDA+"10_stratified_models/shap_houses_rf_bar.png",["Distance to <strong>Cebu Business Park</strong> is the single strongest driver.","Classic bid-rent: detached-house value falls with distance from the core.","Amenities enter only as the bundled MCRAI composite."])
insight("What drives price — vacant lots (RQ1)",EDA+"10_stratified_models/shap_lot_rf_bar.png",["Distance to Cebu Business Park leads — the steepest bid-rent gradient of the three.","Individual amenities are the second-largest driver — bare land responds to specific nearby access.","Still the weakest stratum: unrecorded parcel attributes (frontage, zoning) cap accuracy."])
insight("The valuation gap (RQ4)",DECK+"valuation_gap_lots.png",["Market prices run ~<strong>3×</strong> the BIR zonal benchmark for vacant lots.","Systematic, not noise — listings beat BIR in nearly every LGU.","A signal that benchmarks are stale — not a correction factor to apply."])
points("Is it good enough?",["A triangulation reference, not a replacement.","Strongest for condos and houses; indicative for lots.","It complements professional judgment — never replaces it."])
divider("09","The Walkthrough")
plot("Market Map",DECK+"webapp_market_map.png","Open-market listings with filters.")
plot("Price Surface",DECK+"webapp_price_surface.png","Predicted ₱/sqm by barangay.")
plot("Property Predictor",DECK+"webapp_predictor.png","A live estimate — and its SHAP reasoning.")
divider("10","The Close")
points("Answers to the four questions",["<strong>RQ1</strong> — Location dominates; drivers differ by type.","<strong>RQ2</strong> — Random Forest deployed.","<strong>RQ3</strong> — Geospatial helps, most where benchmarks are weak.","<strong>RQ4</strong> — The BIR gap is large and systematic."])
points("Contribution",["<strong>Methodological</strong> — per-type modeling, polycentric distances, the two-stage MCRAI, leak-free CV.","<strong>Practical</strong> — a reproducible, explainable web prototype."])
points("Limitations",["Asking-price ceiling (not deed-of-sale).","Cross-sectional snapshot.","Vacant-lot data ceiling."])
points("Recommendations — practice & policy",["Use as triangulation, not a final number.","Recognize secondary subcenters; MCRAI as a template, not a finished index."])
points("Recommendations — future research",["Estimate MCRAI parameters from data.","Enrich parcel attributes; add a time dimension."],c="Udomsap & Abid, 2020")
close("A clearer starting point for the family deciding what a home is worth —","not the last word on it.")
references([
    "Agosto, A. (2020). Determinants of Land Values in Cebu City, Philippines.",
    "Bangko Sentral ng Pilipinas. (2025). Residential Real Estate Price Index (RREPI): Q2 2025 Report.",
    "Boeing, G. (2017). OSMnx: New methods for acquiring, constructing, analyzing, and visualizing complex street networks.",
    "Breiman, L. (2001). Random forests.",
    "Bureau of Internal Revenue. (n.d.). Zonal Values Resources and Schedules.",
    "Chen, T., & Guestrin, C. (2016). XGBoost: A scalable tree boosting system. ACM SIGKDD.",
    "Dröes, M. I., Hoesli, M., & Bourassa, S. C. (2019). Heterogeneous households and market segmentation in a hedonic framework. ERES.",
    "Giuliano, G., & Small, K. A. (1991). Subcenters in the Los Angeles region.",
    "Grinsztajn, L., Oyallon, E., & Varoquaux, G. (2022). Why do tree-based models still outperform deep learning on typical tabular data? NeurIPS 35.",
    "Hansen, W. G. (1959). How accessibility shapes land use.",
    "International Valuation Standards Council. (2025). International Valuation Standards (IVS) 2025.",
    "Japan International Cooperation Agency & MCDCB. (2015). The Roadmap Study for Sustainable Urban Development in Metro Cebu.",
    "Lundberg, S. M., & Lee, S.-I. (2017). A unified approach to interpreting model predictions. NeurIPS 30.",
    "McMillen, D. P. (2003). The return of centralization to Chicago.",
    "Otsuka, K., Manasan, R. G., & Piza, S. (2023). Local government unit income and real property tax collection in the Philippines. PIDS.",
    "Philippine Statistics Authority – Central Visayas. (2025). Cebu's economy grows by 7.3 percent in 2024.",
    "Ramolete, G. I. L., Bramaskara, B., Reyes, D. A., & Heinrich, A. (2023). Utilization of machine learning and government-based indicators for property value prediction in the Philippines. The Philippine Statistician.",
    "Rosen, S. (1974). Hedonic prices and implicit markets.",
    "Tanamal, R., Minoque, N., Wiradinata, T., Soekamto, Y., & Ratih, T. (2023). House price prediction model using Random Forest in Surabaya City. TEM Journal.",
    "Tobler, W. R. (1970). A computer movie simulating urban growth in the Detroit region.",
    "Udomsap, A., & Abid, M. (2020). Macroeconomic determinants of housing prices.",
    "Usman, H., Lizam, M., & Adekunle, M. U. (2020). A priori spatial segmentation of commercial property market using hedonic price modelling.",
    "Viray, F. S. (2023). Residential property price forecasting model for Central Pangasinan, Philippines.",
])
thankyou()

out = R("defense_deck.pptx"); prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
