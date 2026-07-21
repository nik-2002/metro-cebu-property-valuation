"""Build defense_deck_v2.pptx from the v2 deck content (build_deck_v2.SLIDES).

Text stays editable; graphs and the MCRAI formula are inserted as images
(PowerPoint cannot reliably typeset the sigma/subscripts, so the formula is a PNG).
Run:  python3 build_pptx_v2.py
"""
import os, re, math, html as _html
from html.parser import HTMLParser
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import build_deck_v2 as deck  # importing builds SLIDES (and harmlessly rewrites the HTML)

HERE = os.path.dirname(os.path.abspath(__file__))
def R(p): return p if os.path.isabs(p) else os.path.join(HERE, p)

# ---- cool palette (matches the v2 HTML) -----------------------------------
BG     = RGBColor(0xFA,0xFB,0xFC); INK   = RGBColor(0x1E,0x25,0x30)
MUTED  = RGBColor(0x6E,0x76,0x86); ACCENT= RGBColor(0x33,0x45,0x5E)
ACCENT2= RGBColor(0x5B,0x7A,0x99); ACCENT3=RGBColor(0xC5,0xD2,0xDF)
SOFT   = RGBColor(0xEA,0xEE,0xF3); LINE  = RGBColor(0xDE,0xE3,0xEA)
WHITE  = RGBColor(0xFF,0xFF,0xFF); CITE  = RGBColor(0x97,0x9F,0xAD)
FONT = "Arial"
# body type sizes (>= 16pt per request; tables a touch smaller so columns fit)
SZ_LEAD, SZ_BUL, SZ_SUB = 20, 18, 16
SZ_SUBH, SZ_NOTE, SZ_TAKE = 17, 14, 18
SZ_TBL, SZ_TBLH, SZ_FEAT, SZ_STATLAB = 13.5, 13.5, 14, 18
EW, EH = 13.333, 7.5
X0, TOP, W, BOTTOM = 0.62, 1.60, 12.10, 7.05

prs = Presentation(); prs.slide_width = Inches(EW); prs.slide_height = Inches(EH)
BLANK = prs.slide_layouts[6]

# ---- mini DOM --------------------------------------------------------------
class Node:
    __slots__ = ("tag", "attrs", "kids")
    def __init__(self, tag=None, attrs=None):
        self.tag = tag; self.attrs = dict(attrs or {}); self.kids = []
    def cls(self): return self.attrs.get("class", "")

class DOM(HTMLParser):
    VOID = {"img", "br"}
    def __init__(self):
        super().__init__(convert_charrefs=True)
        self.root = Node("root"); self.stack = [self.root]
    def handle_starttag(self, t, a):
        n = Node(t, a); self.stack[-1].kids.append(n)
        if t not in self.VOID: self.stack.append(n)
    def handle_startendtag(self, t, a):
        self.stack[-1].kids.append(Node(t, a))
    def handle_endtag(self, t):
        if t in self.VOID: return
        for i in range(len(self.stack) - 1, 0, -1):
            if self.stack[i].tag == t:
                del self.stack[i:]; return
    def handle_data(self, d):
        if d: self.stack[-1].kids.append(d)

def parse(body):
    p = DOM(); p.feed(_html.unescape(body)); return p.root.kids

def node_text(n):
    if isinstance(n, str): return n
    return "".join(node_text(k) for k in n.kids)

# ---- low-level shape helpers ----------------------------------------------
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s

def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    return tf

def rect(s, l, t, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def add_runs(p, node, bold=False, italic=False, color=INK, size=None):
    """Append formatted runs from a DOM node (or string) into paragraph p."""
    kids = [node] if isinstance(node, str) else node.kids
    for k in kids:
        if isinstance(k, str):
            if k == "": continue
            r = p.add_run(); r.text = k
            f = r.font; f.name = FONT; f.bold = bold; f.italic = italic; f.color.rgb = color
            if size: f.size = Pt(size)
        elif k.tag in ("b", "strong"):
            add_runs(p, k, True, italic, color, size)
        elif k.tag == "em":
            add_runs(p, k, bold, True, color, size)
        elif k.tag == "span" and "c" in k.cls().split():
            add_runs(p, k, bold, True, CITE, size)
        elif k.tag == "br":
            r = p.add_run(); r.text = "\n"
        else:
            add_runs(p, k, bold, italic, color, size)

def est_h(text, size, width, lh=1.2, pad=0.07):
    cpl = max(6, int(width * 138 / size))
    lines = max(1, math.ceil(len(text) / cpl))
    return lines * size * lh / 72.0 + pad

def img_fit(s, path, bx, by, bw, bh):
    iw, ih = Image.open(R(path)).size
    sc = min(bw / iw, bh / ih); w, h = iw * sc, ih * sc
    s.shapes.add_picture(R(path), Inches(bx + (bw - w) / 2), Inches(by + (bh - h) / 2),
                         Inches(w), Inches(h))

# ---- block renderers (return height consumed) ------------------------------
def r_lead(s, n, x, y, w):
    h = est_h(node_text(n), SZ_LEAD, w, pad=0.16)
    tf = box(s, x, y, w, h); p = tf.paragraphs[0]; p.line_spacing = 1.18
    add_runs(p, n, color=INK, size=SZ_LEAD)
    return h

def r_subh(s, n, x, y, w):
    tf = box(s, x, y + 0.04, w, 0.38); p = tf.paragraphs[0]
    add_runs(p, n, bold=True, color=ACCENT, size=SZ_SUBH)
    return 0.48

def r_note(s, n, x, y, w, color=MUTED, size=SZ_NOTE):
    h = est_h(node_text(n), size, w, pad=0.08)
    tf = box(s, x, y, w, h); p = tf.paragraphs[0]; p.line_spacing = 1.12
    add_runs(p, n, italic=True, color=color, size=size)
    return h

def r_bullets(s, ul, x, y, w):
    tf = box(s, x, y, w, BOTTOM - y); first = True; total = 0.0
    for li in ul.kids:
        if isinstance(li, str) or li.tag != "li": continue
        # split off any nested sub-list
        subs = [k for k in li.kids if not isinstance(k, str) and k.tag == "ul"]
        main = Node("li"); main.kids = [k for k in li.kids
                                        if isinstance(k, str) or k.tag != "ul"]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.line_spacing = 1.18; p.space_after = Pt(8)
        rb = p.add_run(); rb.text = "▪  "; rb.font.name = FONT; rb.font.size = Pt(14)
        rb.font.color.rgb = ACCENT; rb.font.bold = True
        add_runs(p, main, color=INK, size=SZ_BUL)
        total += est_h(node_text(main), SZ_BUL, w - 0.3, pad=0.14)
        for sub in subs:
            for sli in sub.kids:
                if isinstance(sli, str) or sli.tag != "li": continue
                sp = tf.add_paragraph(); sp.line_spacing = 1.14; sp.space_after = Pt(4)
                sp.level = 1
                rb = sp.add_run(); rb.text = "–  "; rb.font.name = FONT; rb.font.size = Pt(14)
                rb.font.color.rgb = ACCENT2
                add_runs(sp, sli, color=MUTED, size=SZ_SUB)
                total += est_h(node_text(sli), SZ_SUB, w - 0.7, pad=0.10)
    return total + 0.08

def r_takeaway(s, n, x, y, w):
    txt = "".join(node_text(k) for k in n.kids if not (isinstance(k, Node) and k.tag == "span"))
    h = max(0.7, est_h(txt, SZ_TAKE, w - 0.9, pad=0.36))
    rect(s, x, y, w, h, SOFT)
    rect(s, x, y, 0.07, h, ACCENT)
    tf = box(s, x + 0.32, y, w - 0.5, h, MSO_ANCHOR.MIDDLE); p = tf.paragraphs[0]
    p.line_spacing = 1.16
    rk = p.add_run(); rk.text = "TAKEAWAY   "; rk.font.name = FONT; rk.font.bold = True
    rk.font.size = Pt(12); rk.font.color.rgb = ACCENT2
    # the takeaway helper wraps label in a <span>; skip it, render the trailing text
    tail = Node("x"); tail.kids = [k for k in n.kids if not (isinstance(k, Node) and k.tag == "span")]
    add_runs(p, tail, color=INK, size=SZ_TAKE)
    return h + 0.18

def r_statrow(s, n, x, y, w):
    big = next((k for k in n.kids if isinstance(k, Node) and "bigstat" in k.cls()), None)
    lab = next((k for k in n.kids if isinstance(k, Node) and "statlab" in k.cls()), None)
    tf = box(s, x, y, w, 1.2)
    p = tf.paragraphs[0]
    rb = p.add_run(); rb.text = node_text(big) if big else ""
    rb.font.name = FONT; rb.font.bold = True; rb.font.size = Pt(46); rb.font.color.rgb = ACCENT
    if lab:
        p2 = tf.add_paragraph(); p2.line_spacing = 1.15; p2.space_before = Pt(2)
        add_runs(p2, lab, color=INK, size=SZ_STATLAB)
    return 1.3

def r_table(s, tbl, x, y, w, note=None):
    headers = []; rows = []
    thead = next((k for k in tbl.kids if isinstance(k, Node) and k.tag == "thead"), None)
    tbody = next((k for k in tbl.kids if isinstance(k, Node) and k.tag == "tbody"), None)
    if thead:
        tr = next(k for k in thead.kids if isinstance(k, Node) and k.tag == "tr")
        headers = [c for c in tr.kids if isinstance(c, Node) and c.tag == "th"]
    if tbody:
        for tr in tbody.kids:
            if isinstance(tr, str) or tr.tag != "tr": continue
            rows.append([c for c in tr.kids if isinstance(c, Node) and c.tag == "td"])
    nr, nc = len(rows) + 1, len(headers)
    rowh = 0.46
    h = rowh * nr
    gfx = s.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h))
    t = gfx.table
    t.first_row = False
    # widths: first column wider
    t.columns[0].width = Inches(w * (0.34 if nc <= 3 else 0.24))
    rest = (w - (w * (0.34 if nc <= 3 else 0.24))) / max(1, nc - 1)
    for j in range(1, nc): t.columns[j].width = Inches(rest)
    for j, th in enumerate(headers):
        cell = t.cell(0, j); cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        cell.margin_top = cell.margin_bottom = Pt(3)
        p = cell.text_frame.paragraphs[0]
        add_runs(p, th, bold=True, color=WHITE, size=SZ_TBLH)
    for i, row in enumerate(rows, 1):
        for j in range(nc):
            cell = t.cell(i, j); cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else SOFT
            cell.margin_top = cell.margin_bottom = Pt(3)
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            if j < len(row): add_runs(p, row[j], color=INK, size=SZ_TBL)
    used = h + 0.1
    if note is not None:
        used += r_note(s, note, x, y + h + 0.06, w)
    return used + 0.06

def r_formula(s, x, y, w):
    img_fit(s, DECK + "formula_mcrai.png", x, y, min(w, 6.0), 0.95)
    return 1.05

def r_blocks(s, kids, x, y, w, bottom):
    cy = y; i = 0
    while i < len(kids):
        n = kids[i]; i += 1
        if isinstance(n, str):
            continue
        c = n.cls()
        if n.tag == "p" and "lead" in c:           cy += r_lead(s, n, x, cy, w)
        elif n.tag == "p" and "note-inline" in c:  cy += r_note(s, n, x, cy, w)
        elif n.tag == "div" and "subh" in c:        cy += r_subh(s, n, x, cy, w)
        elif n.tag == "ul":                          cy += r_bullets(s, n, x, cy, w)
        elif n.tag == "div" and "takeaway" in c:     cy += r_takeaway(s, n, x, cy, w)
        elif n.tag == "div" and "statrow" in c:      cy += r_statrow(s, n, x, cy, w)
        elif n.tag == "div" and "formula" in c:      cy += r_formula(s, x, cy, w)
        elif n.tag == "table":
            note = None
            if i < len(kids) and isinstance(kids[i], Node) and "tnote" in kids[i].cls():
                note = kids[i]; i += 1
            cy += r_table(s, n, x, cy, w, note)
        elif n.tag == "div" and "twocol" in c:       cy += r_twocol(s, n, x, cy, w, bottom)
        elif n.tag == "div" and "featgrid" in c:     cy += r_featgrid(s, n, x, cy, w, bottom)
        elif n.tag == "div" and c.split()[:1] == ["fa"]: cy += r_fa(s, n, x, cy, w, bottom)
        elif n.tag == "div" and "fig-full" in c:     cy += r_figfull(s, n, x, cy, w, bottom)
    return cy - y

def r_twocol(s, n, x, y, w, bottom):
    cols = [k for k in n.kids if isinstance(k, Node) and k.tag == "div"]
    gap = 0.5; cw = (w - gap) / 2
    h1 = r_blocks(s, cols[0].kids, x, y, cw, bottom) if cols else 0
    h2 = r_blocks(s, cols[1].kids, x + cw + gap, y, cw, bottom) if len(cols) > 1 else 0
    return max(h1, h2) + 0.1

def r_featgrid(s, n, x, y, w, bottom):
    cols = [k for k in n.kids if isinstance(k, Node) and "fcol" in k.cls()]
    if not cols: return 0
    gap = 0.35; cw = (w - gap * (len(cols) - 1)) / len(cols); hmax = 0
    for j, col in enumerate(cols):
        cx = x + j * (cw + gap); cy = y
        fh = next((k for k in col.kids if isinstance(k, Node) and "fh" in k.cls()), None)
        if fh:
            tf = box(s, cx, cy, cw, 0.36); add_runs(tf.paragraphs[0], fh, bold=True, color=ACCENT, size=14)
            cy += 0.44
        ul = next((k for k in col.kids if isinstance(k, Node) and k.tag == "ul"), None)
        if ul:
            tf = box(s, cx, cy, cw, bottom - cy); first = True
            for li in ul.kids:
                if isinstance(li, str) or li.tag != "li": continue
                p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
                p.line_spacing = 1.14; p.space_after = Pt(4)
                add_runs(p, li, color=INK, size=SZ_FEAT)
                cy += est_h(node_text(li), SZ_FEAT, cw, pad=0.09)
        hmax = max(hmax, cy - y)
    return hmax + 0.1

def r_fa(s, n, x, y, w, bottom):
    fign = next((k for k in n.kids if isinstance(k, Node) and "fa-fig" in k.cls()), None)
    txtn = next((k for k in n.kids if isinstance(k, Node) and "fa-txt" in k.cls()), None)
    img = next((k for k in fign.kids if isinstance(k, Node) and k.tag == "img"), None)
    src = img.attrs.get("src") if img else None
    wide = "wide" in (fign.cls().split() if fign else [])
    avail = bottom - y
    if wide:
        ih = avail * 0.60
        if src: img_fit(s, src, x, y, w, ih)
        if txtn: r_blocks(s, txtn.kids, x, y + ih + 0.1, w, bottom)
    else:
        iw = w * 0.54
        if src: img_fit(s, src, x, y, iw, avail)
        if txtn: r_blocks(s, txtn.kids, x + iw + 0.3, y, w - iw - 0.3, bottom)
    return avail

def r_figfull(s, n, x, y, w, bottom):
    img = next((k for k in n.kids if isinstance(k, Node) and k.tag == "img"), None)
    cap = next((k for k in n.kids if isinstance(k, Node) and "figcap" in k.cls()), None)
    tk  = next((k for k in n.kids if isinstance(k, Node) and "takeaway" in k.cls()), None)
    avail = bottom - y; reserve = (0.6 if cap else 0) + (0.8 if tk else 0)
    ih = avail - reserve
    if img: img_fit(s, img.attrs.get("src"), x, y, w, ih)
    cy = y + ih + 0.05
    if cap:
        tf = box(s, x, cy, w, 0.55); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        add_runs(p, cap, italic=True, color=MUTED, size=13); cy += 0.6
    if tk: r_takeaway(s, tk, x, cy, w)
    return avail

# ---- measurement pass (for vertical centring of text slides) ---------------
def est_one(n, w):
    c = n.cls(); t = n.tag
    if t == "p" and "lead" in c:        return est_h(node_text(n), SZ_LEAD, w, pad=0.16)
    if t == "p" and "note-inline" in c: return est_h(node_text(n), SZ_NOTE, w, pad=0.08)
    if t == "div" and "subh" in c:      return 0.48
    if t == "div" and "statrow" in c:   return 1.30
    if t == "div" and "formula" in c:   return 1.05
    if t == "ul":
        tot = 0.0
        for li in n.kids:
            if isinstance(li, str) or li.tag != "li": continue
            subs = [k for k in li.kids if isinstance(k, Node) and k.tag == "ul"]
            mn = Node("li"); mn.kids = [k for k in li.kids if isinstance(k, str) or k.tag != "ul"]
            tot += est_h(node_text(mn), SZ_BUL, w - 0.3, pad=0.14)
            for sub in subs:
                for sli in sub.kids:
                    if isinstance(sli, str) or sli.tag != "li": continue
                    tot += est_h(node_text(sli), SZ_SUB, w - 0.7, pad=0.10)
        return tot + 0.08
    if t == "div" and "takeaway" in c:
        txt = "".join(node_text(k) for k in n.kids if not (isinstance(k, Node) and k.tag == "span"))
        return max(0.7, est_h(txt, SZ_TAKE, w - 0.9, pad=0.36)) + 0.18
    return None  # fa / fig-full -> a "fill" slide, do not centre

def est_blocks(kids, w):
    tot = 0.0; i = 0
    while i < len(kids):
        n = kids[i]; i += 1
        if isinstance(n, str): continue
        c = n.cls()
        if n.tag == "table":
            nr = 1 + sum(1 for k in (next((x for x in n.kids if isinstance(x, Node) and x.tag == "tbody"), Node("x")).kids)
                         if isinstance(k, Node) and k.tag == "tr")
            h = 0.46 * nr + 0.1
            if i < len(kids) and isinstance(kids[i], Node) and "tnote" in kids[i].cls():
                h += est_h(node_text(kids[i]), SZ_NOTE, w, pad=0.08); i += 1
            tot += h + 0.06
        elif n.tag == "div" and "twocol" in c:
            cols = [k for k in n.kids if isinstance(k, Node) and k.tag == "div"]; cw = (w - 0.5) / 2
            h1 = est_blocks(cols[0].kids, cw) if cols else 0
            h2 = est_blocks(cols[1].kids, cw) if len(cols) > 1 else 0
            if h1 is None or h2 is None: return None
            tot += max(h1, h2) + 0.1
        elif n.tag == "div" and "featgrid" in c:
            cols = [k for k in n.kids if isinstance(k, Node) and "fcol" in k.cls()]
            cw = (w - 0.35 * (len(cols) - 1)) / max(1, len(cols)); hm = 0
            for col in cols:
                ch = 0.44 if any(isinstance(k, Node) and "fh" in k.cls() for k in col.kids) else 0
                ul = next((k for k in col.kids if isinstance(k, Node) and k.tag == "ul"), None)
                if ul:
                    for li in ul.kids:
                        if isinstance(li, Node) and li.tag == "li":
                            ch += est_h(node_text(li), SZ_FEAT, cw, pad=0.09)
                hm = max(hm, ch)
            tot += hm + 0.1
        else:
            h = est_one(n, w)
            if h is None: return None
            tot += h
    return tot

# ---- slide-level builders --------------------------------------------------
DECK = deck.DECK if hasattr(deck, "DECK") else "assets/deck/"

def header(s, section, title, page, n):
    tf = box(s, X0, 0.46, 9.0, 0.4)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = _html.unescape(section).upper()
    r.font.name = FONT; r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = ACCENT
    pg = box(s, EW - 2.2, 0.46, 1.55, 0.4)
    pr = pg.paragraphs[0]; pr.alignment = PP_ALIGN.RIGHT
    r = pr.add_run(); r.text = f"{page} / {n}"; r.font.name = FONT; r.font.size = Pt(12.5)
    r.font.color.rgb = MUTED
    rect(s, X0, 0.96, EW - 2 * X0, 0.018, ACCENT)
    tt = box(s, X0, 1.04, EW - 2 * X0, 0.55)
    add_runs(tt.paragraphs[0], parse(title)[0] if title else Node("x"), bold=True, color=INK, size=30)

def title_slide(s):
    tf = box(s, X0, 2.0, EW - 2 * X0, 0.5)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "University of Asia and the Pacific · BS Data Science · June 2026"
    r.font.name = FONT; r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = ACCENT
    t2 = box(s, X0, 2.6, EW - 2 * X0, 1.8)
    p = t2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.line_spacing = 1.08
    r = p.add_run(); r.text = "Predicting Open-Market Residential Property Values in Metro Cebu"
    r.font.name = FONT; r.font.bold = True; r.font.size = Pt(38); r.font.color.rgb = INK
    rect(s, EW / 2 - 0.5, 4.55, 1.0, 0.05, ACCENT)
    t3 = box(s, X0, 4.8, EW - 2 * X0, 1.4)
    p = t3.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Using Machine Learning and Geospatial Features"
    r.font.name = FONT; r.font.size = Pt(21); r.font.color.rgb = MUTED
    p2 = t3.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(18)
    r = p2.add_run(); r.text = "Chris Dominic Estreba"; r.font.name = FONT; r.font.size = Pt(18)
    r.font.color.rgb = INK
    p3 = t3.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r = p3.add_run(); r.text = "BS Data Science Capstone"; r.font.name = FONT; r.font.size = Pt(15)
    r.font.color.rgb = MUTED

def divider_slide(s, num, title, sub, dark):
    if dark:
        tf = box(s, X0, 2.6, EW - 2 * X0, 2.3, MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = _html.unescape(title); r.font.name = FONT; r.font.bold = True
        r.font.size = Pt(60); r.font.color.rgb = WHITE
        if sub:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(10)
            r = p2.add_run(); r.text = _html.unescape(sub); r.font.name = FONT
            r.font.size = Pt(20); r.font.color.rgb = ACCENT3
        return
    tf = box(s, X0 + 0.3, 2.4, EW - 2 * X0, 2.6, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = _html.unescape(num); r.font.name = FONT; r.font.bold = True
    r.font.size = Pt(40); r.font.color.rgb = ACCENT2
    p2 = tf.add_paragraph(); p2.space_before = Pt(6)
    r = p2.add_run(); r.text = _html.unescape(title); r.font.name = FONT; r.font.bold = True
    r.font.size = Pt(50); r.font.color.rgb = WHITE
    if sub:
        p3 = tf.add_paragraph(); p3.space_before = Pt(8)
        r = p3.add_run(); r.text = _html.unescape(sub); r.font.name = FONT
        r.font.size = Pt(19); r.font.color.rgb = ACCENT3

# ---- main ------------------------------------------------------------------
def main():
    SL = deck.SLIDES; n = len(SL)
    for i, sd in enumerate(SL, 1):
        if i == 1:
            title_slide(slide()); continue
        if sd["dnum"]:
            divider_slide(slide(ACCENT), sd["dnum"], sd["title"], sd["dsub"], False); continue
        if sd["dark"]:
            divider_slide(slide(ACCENT), "", sd["title"], sd["dsub"], True); continue
        s = slide()
        header(s, sd["section"], sd["title"], i, n)
        kids = parse(sd["body"])
        H = est_blocks(kids, W)
        if H is None:                       # figure slide -> fills the area, keep at top
            start = TOP
        else:                               # text slide -> centre in the body band (bottom ~2/3)
            start = TOP + max(0.0, (BOTTOM - TOP - H) * 0.5)
        r_blocks(s, kids, X0, start, W, BOTTOM)
    out = os.path.join(HERE, "defense_deck_v2.pptx")
    prs.save(out)
    print(f"wrote defense_deck_v2.pptx  ({n} slides)")

if __name__ == "__main__":
    main()
