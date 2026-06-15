"""Append a 'Conclusion Synthesis' slide with copy-pasteable closing text."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent.parent
SRC = ROOT / "final_Draft_Estreba_thesis_proposal.pptx"

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
AMBER = RGBColor(0xD4, 0xA8, 0x43)
SLATE = RGBColor(0xEE, 0xF1, 0xF6)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x5F, 0x6D)


def add_header(slide, slide_w, eyebrow, title):
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, Inches(0.55))
    strip.line.fill.background()
    strip.fill.solid(); strip.fill.fore_color.rgb = NAVY
    tf = strip.text_frame
    tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = eyebrow
    p.font.name = "Inter"; p.font.size = Pt(11); p.font.bold = True
    p.font.color.rgb = AMBER

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7),
                                         slide_w - Inches(1.0), Inches(0.7))
    tf = title_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Inter"; p.font.size = Pt(28); p.font.bold = True
    p.font.color.rgb = NAVY

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(1.45), Inches(0.9), Inches(0.05))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = AMBER


def main():
    prs = Presentation(str(SRC))
    sw = prs.slide_width
    s = prs.slides.add_slide(prs.slide_layouts[0])

    add_header(s, sw,
               "CONCLUSION  /  SYNTHESIS",
               "What This Thesis Proves")

    # Centered hero quote — large, copy-pasteable
    quote_left = Inches(1.0)
    quote_top = Inches(2.0)
    quote_w = sw - Inches(2.0)
    quote_h = Inches(3.5)

    qb = s.shapes.add_textbox(quote_left, quote_top, quote_w, quote_h)
    tf = qb.text_frame; tf.word_wrap = True

    # Sentence 1
    p = tf.paragraphs[0]
    p.text = ("This thesis proves that property-level, geospatially grounded "
              "valuation is feasible in Metro Cebu using only publicly "
              "available data.")
    p.font.name = "Inter"; p.font.size = Pt(22); p.font.color.rgb = NAVY
    p.font.bold = True
    p.space_after = Pt(18)

    # Sentence 2
    p = tf.add_paragraph()
    p.text = ("The 59% MAPE is the honest ceiling of asking-price data.")
    p.font.name = "Inter"; p.font.size = Pt(18); p.font.color.rgb = DARK
    p.space_after = Pt(8)

    # Sentence 3
    p = tf.add_paragraph()
    p.text = ("The per-sqm MAE of one BIR zonal tier is the honest floor "
              "of practical utility.")
    p.font.name = "Inter"; p.font.size = Pt(18); p.font.color.rgb = DARK
    p.space_after = Pt(18)

    # Closing
    p = tf.add_paragraph()
    p.text = ("Defensible, reproducible, and ready for practitioners today.")
    p.font.name = "Inter"; p.font.size = Pt(18); p.font.italic = True
    p.font.color.rgb = AMBER
    p.font.bold = True

    # Bottom amber accent
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(1.0), Inches(6.4),
                                Inches(0.9), Inches(0.06))
    accent.line.fill.background()
    accent.fill.solid(); accent.fill.fore_color.rgb = AMBER

    # Footer note
    footer = s.shapes.add_textbox(Inches(1.0), Inches(6.55),
                                  sw - Inches(2.0), Inches(0.4))
    p = footer.text_frame.paragraphs[0]
    p.text = "Copy-pasteable synthesis  ·  use as the lead paragraph on the closing slide."
    p.font.name = "Inter"; p.font.size = Pt(10); p.font.italic = True
    p.font.color.rgb = GRAY

    prs.save(str(SRC))
    print(f"Saved: {SRC.name}  ·  total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
